from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from typing import Dict, Any
import pandas as pd

class OfficeLoader:
    def load_docx(self, file_path: str) -> Dict[str, Any]:
        """Load Word document"""
        doc = Document(file_path)
        
        paragraphs = []
        tables = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append({
                    "text": para.text,
                    "style": para.style.name if para.style else None
                })
        
        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(pd.DataFrame(table_data))
        
        return {
            "paragraphs": paragraphs,
            "tables": tables,
            "full_text": "\n\n".join([p["text"] for p in paragraphs]),
            "metadata": {
                "num_paragraphs": len(paragraphs),
                "num_tables": len(tables)
            }
        }
    
    def load_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Load Excel file"""
        wb = load_workbook(file_path, data_only=True)
        
        sheets = {}
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            data = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    data.append(list(row))
            
            if data:
                df = pd.DataFrame(data[1:], columns=data[0] if data else [])
                sheets[sheet_name] = df
        
        # Create text representation
        full_text = ""
        for sheet_name, df in sheets.items():
            full_text += f"\n## Sheet: {sheet_name}\n"
            full_text += df.to_string()
        
        return {
            "sheets": sheets,
            "full_text": full_text,
            "metadata": {
                "num_sheets": len(sheets),
                "sheet_names": list(sheets.keys())
            }
        }
    
    def load_pptx(self, file_path: str) -> Dict[str, Any]:
        """Load PowerPoint presentation"""
        prs = Presentation(file_path)
        
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            slides.append({
                "slide_num": i + 1,
                "text": "\n".join(slide_text),
                "layout": slide.slide_layout.name if slide.slide_layout else None
            })
        
        return {
            "slides": slides,
            "full_text": "\n\n---\n\n".join([f"Slide {s['slide_num']}:\n{s['text']}" for s in slides]),
            "metadata": {
                "num_slides": len(slides)
            }
        }